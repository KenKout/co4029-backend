"""PPTX extractor — python-pptx wrapping with slide-level source locations."""

from __future__ import annotations

import asyncio
import io
from typing import BinaryIO

from abridgeai.ai.extraction.base import ExtractedContent, SourceLocation
from abridgeai.ai.extraction.registry import register_extractor

_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


@register_extractor(_PPTX_MIME)
class PptxExtractor:
    async def extract(self, source: BinaryIO | bytes | str) -> ExtractedContent:
        return await asyncio.to_thread(_extract_sync, source)


def _extract_sync(source: BinaryIO | bytes | str) -> ExtractedContent:
    from pptx import Presentation

    handle = io.BytesIO(source) if isinstance(source, bytes) else source
    presentation = Presentation(handle)

    parts: list[str] = []
    locations: list[SourceLocation] = []
    slide_total = 0
    for slide_num, slide in enumerate(presentation.slides, start=1):
        slide_total = slide_num
        slide_lines = [
            getattr(shape, "text", "").strip()
            for shape in slide.shapes
            if getattr(shape, "text", "").strip()
        ]
        if slide_lines:
            parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_lines))
            locations.append(SourceLocation(page=slide_num))

    body = "\n\n".join(parts).strip()
    return ExtractedContent(
        text=body,
        metadata={"slide_count": slide_total},
        source_type="pptx",
        source_locations=locations,
    )


__all__ = ["PptxExtractor"]
