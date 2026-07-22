"""Generate synthetic multimodal test fixtures (T4.8).

Creates 7 small, license-clean fixture files used by the extractor test
suite. All content is original/synthetic - see ``LICENSE.md``. Total
size stays under 500KB to satisfy the pre-commit ``check-added-large-files``
hook (T0.22).

Usage::

    cd backend-new/tests/fixtures/multimodal
    uv run python generate_fixtures.py

The script is idempotent: it overwrites existing fixtures on each run.
Whenever a fixture is regenerated, re-run the extractor tests
(``uv run pytest tests/unit/test_extractors_with_fixtures.py``) to confirm
assertions still hold.
"""

from __future__ import annotations

import datetime as _dt
import io
import math
import re
import shutil
import struct
import subprocess
import sys
import wave
import zipfile
from pathlib import Path

OUT_DIR = Path(__file__).resolve().parent

FIXED_TIMESTAMP = _dt.datetime(2024, 1, 1, 0, 0, 0, tzinfo=_dt.UTC)
FIXED_ZIP_DATE_TIME = (2024, 1, 1, 0, 0, 0)
PDF_FIXED_DATE = "D:20240101000000Z"
_PDF_ID_PATTERN = re.compile(rb"/ID\s*\[[^\]]+\]")
_PDF_FIXED_ID_BYTES = b"/ID[<00000000000000000000000000000001><00000000000000000000000000000001>]"


def _normalize_zip(path: Path) -> None:
    with path.open("rb") as fh:
        original = fh.read()
    buf = io.BytesIO()
    with zipfile.ZipFile(io.BytesIO(original), "r") as src:
        names = sorted(src.namelist())
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as dst:
            for name in names:
                data = src.read(name)
                info = zipfile.ZipInfo(filename=name, date_time=FIXED_ZIP_DATE_TIME)
                info.compress_type = zipfile.ZIP_DEFLATED
                info.create_system = 0
                info.external_attr = 0
                dst.writestr(info, data)
    path.write_bytes(buf.getvalue())


def _normalize_pdf_id(path: Path) -> None:
    data = path.read_bytes()
    replaced = _PDF_ID_PATTERN.sub(_PDF_FIXED_ID_BYTES, data, count=1)
    if replaced != data:
        path.write_bytes(replaced)


PDF_FILENAME = "sample.pdf"
DOCX_FILENAME = "sample.docx"
PPTX_FILENAME = "sample.pptx"
WAV_FILENAME = "sample.wav"
MP4_FILENAME = "sample.mp4"
PNG_FILENAME = "text-image.png"
HTML_FILENAME = "sample.html"
XLSX_FILENAME = "sample.xlsx"

WAV_SAMPLE_RATE_HZ = 16000
WAV_DURATION_S = 5
WAV_FREQUENCY_HZ = 440
WAV_BIT_DEPTH_BYTES = 2
WAV_AMPLITUDE = 0.3
PCM_INT16_MAX = 32767


def _gen_pdf(out: Path) -> None:
    import fitz  # type: ignore[import-untyped,unused-ignore]

    doc = fitz.open()
    page1 = doc.new_page()
    page1.insert_text((72, 72), "Hello World")
    page1.insert_text((72, 200), "This is page 1 of the test fixture.")
    page2 = doc.new_page()
    page2.insert_text((72, 72), "Page 2 content")
    page2.insert_text((72, 200), "Second page for multi-page extraction tests.")
    doc.set_metadata(
        {
            "title": "Multimodal Test Fixture",
            "author": "abridgeai test suite",
            "subject": "synthetic",
            "creator": "generate_fixtures.py",
            "producer": "generate_fixtures.py",
            "creationDate": PDF_FIXED_DATE,
            "modDate": PDF_FIXED_DATE,
        }
    )
    doc.save(str(out))
    doc.close()
    _normalize_pdf_id(out)


def _gen_docx(out: Path) -> None:
    from docx import Document

    d = Document()
    d.add_heading("Test Fixture", level=1)
    d.add_paragraph("Hello World - this is a synthetic DOCX fixture for extraction tests.")
    core_props = d.core_properties
    core_props.created = FIXED_TIMESTAMP
    core_props.modified = FIXED_TIMESTAMP
    core_props.last_modified_by = ""
    core_props.author = "abridgeai"
    core_props.title = "Test Fixture"
    d.save(str(out))
    _normalize_zip(out)


def _gen_pptx(out: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    p = Presentation()
    title_only_layout = p.slide_layouts[5]

    s1 = p.slides.add_slide(title_only_layout)
    s1.shapes.title.text = "Slide 1: Hello World"
    body1 = s1.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1))
    body1.text_frame.text = "First slide body content."

    s2 = p.slides.add_slide(title_only_layout)
    s2.shapes.title.text = "Slide 2: Test fixture"
    body2 = s2.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1))
    body2.text_frame.text = "Second slide body content."

    core_props = p.core_properties
    core_props.created = FIXED_TIMESTAMP
    core_props.modified = FIXED_TIMESTAMP
    core_props.last_modified_by = ""
    core_props.author = "abridgeai"
    core_props.title = "Test Fixture"

    p.save(str(out))
    _normalize_zip(out)


def _gen_wav(out: Path) -> None:
    n_samples = WAV_SAMPLE_RATE_HZ * WAV_DURATION_S
    frames = bytearray()
    for i in range(n_samples):
        value = int(
            PCM_INT16_MAX
            * WAV_AMPLITUDE
            * math.sin(2 * math.pi * WAV_FREQUENCY_HZ * i / WAV_SAMPLE_RATE_HZ)
        )
        frames += struct.pack("<h", value)

    with wave.open(str(out), "wb") as w:
        w.setnchannels(1)
        w.setsampwidth(WAV_BIT_DEPTH_BYTES)
        w.setframerate(WAV_SAMPLE_RATE_HZ)
        w.writeframes(bytes(frames))


def _gen_mp4(out: Path) -> bool:
    """Returns True on success, False when ffmpeg is unavailable.

    Tests gate on ``out.exists()`` so a False return path leaves the
    suite skipping cleanly rather than failing.
    """
    if not shutil.which("ffmpeg"):
        sys.stderr.write(
            "WARNING: ffmpeg not on PATH; sample.mp4 not generated. "
            "Install ffmpeg and re-run to populate the video fixture.\n"
        )
        return False

    ffmpeg_path = shutil.which("ffmpeg")
    assert ffmpeg_path is not None
    subprocess.run(  # noqa: S603
        [
            ffmpeg_path,
            "-y",
            "-f",
            "lavfi",
            "-i",
            "testsrc=duration=5:size=320x240:rate=10",
            "-f",
            "lavfi",
            "-i",
            "sine=frequency=440:duration=5",
            "-shortest",
            "-pix_fmt",
            "yuv420p",
            str(out),
        ],
        check=True,
        capture_output=True,
    )
    return True


def _gen_png(out: Path) -> None:
    from PIL import Image, ImageDraw

    img = Image.new("RGB", (200, 100), color=(255, 255, 255))
    d = ImageDraw.Draw(img)
    d.text((10, 30), "Hello World", fill=(0, 0, 0))
    img.save(str(out), format="PNG", optimize=True)


def _gen_html(out: Path) -> None:
    out.write_text(
        "<!DOCTYPE html><html><head><title>Test</title></head>"
        "<body><h1>Title</h1><p>Body</p></body></html>\n",
        encoding="utf-8",
    )


def _gen_xlsx(out: Path) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws1 = wb.active
    assert ws1 is not None  # noqa: S101  -- fresh Workbook always has an active sheet
    ws1.title = "Sheet1"
    ws1.append(["Name", "Score"])
    ws1.append(["Hello World", 42])
    ws2 = wb.create_sheet("Sheet2")
    ws2.append(["Second sheet content"])
    wb.properties.created = FIXED_TIMESTAMP
    wb.properties.modified = FIXED_TIMESTAMP
    wb.properties.creator = "abridgeai"
    wb.properties.lastModifiedBy = ""
    wb.properties.title = "Test Fixture"
    wb.save(str(out))
    _normalize_zip(out)


def main() -> int:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    _gen_pdf(OUT_DIR / PDF_FILENAME)
    _gen_docx(OUT_DIR / DOCX_FILENAME)
    _gen_pptx(OUT_DIR / PPTX_FILENAME)
    _gen_wav(OUT_DIR / WAV_FILENAME)
    mp4_ok = _gen_mp4(OUT_DIR / MP4_FILENAME)
    _gen_png(OUT_DIR / PNG_FILENAME)
    _gen_html(OUT_DIR / HTML_FILENAME)
    _gen_xlsx(OUT_DIR / XLSX_FILENAME)

    sys.stdout.write(f"Generated fixtures in {OUT_DIR}\n")
    for name in (
        PDF_FILENAME,
        DOCX_FILENAME,
        PPTX_FILENAME,
        WAV_FILENAME,
        MP4_FILENAME,
        PNG_FILENAME,
        HTML_FILENAME,
        XLSX_FILENAME,
    ):
        path = OUT_DIR / name
        if path.exists():
            sys.stdout.write(f"  {name:20s} {path.stat().st_size:>8d} bytes\n")
        else:
            sys.stdout.write(f"  {name:20s} <skipped>\n")
    if not mp4_ok:
        sys.stdout.write(
            "\nNote: sample.mp4 was skipped (ffmpeg unavailable). "
            "Video extractor tests will skip until ffmpeg is installed.\n"
        )
    return 0


if __name__ == "__main__":
    sys.exit(main())
