from __future__ import annotations

import io

import pytest

from abridgeai.ai.extraction import (
    DocxExtractor,
    ExtractedContent,
    MaterialExtractor,
    PdfExtractor,
    PptxExtractor,
    SourceLocation,
)


@pytest.fixture
def sample_pdf_bytes() -> bytes:
    import fitz

    doc = fitz.open()
    page1 = doc.new_page()
    page1.insert_text((72, 72), "Hello World page one")
    page2 = doc.new_page()
    page2.insert_text((72, 72), "Second page payload")
    data = doc.tobytes()
    doc.close()
    return bytes(data)


@pytest.fixture
def sample_docx_bytes() -> bytes:
    import docx

    document = docx.Document()
    document.add_paragraph("First paragraph hello world")
    document.add_paragraph("")
    document.add_paragraph("Second paragraph extraction test")
    buf = io.BytesIO()
    document.save(buf)
    return buf.getvalue()


@pytest.fixture
def sample_pptx_bytes() -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    slide_layout = presentation.slide_layouts[5]
    slide_one = presentation.slides.add_slide(slide_layout)
    slide_one.shapes.title.text = "Slide One Title"
    slide_two = presentation.slides.add_slide(slide_layout)
    slide_two.shapes.title.text = "Slide Two Title"
    buf = io.BytesIO()
    presentation.save(buf)
    return buf.getvalue()


@pytest.mark.asyncio
async def test_pdf_extract(sample_pdf_bytes: bytes) -> None:
    result = await PdfExtractor().extract(sample_pdf_bytes)

    assert isinstance(result, ExtractedContent)
    assert "Hello World" in result.text
    assert "Second page" in result.text
    assert result.source_type == "pdf"
    assert result.metadata["page_count"] == 2
    assert len(result.source_locations) == 2
    assert result.source_locations[0].page == 1
    assert result.source_locations[1].page == 2


@pytest.mark.asyncio
async def test_pdf_extract_from_file_path(tmp_path, sample_pdf_bytes: bytes) -> None:
    pdf_path = tmp_path / "sample.pdf"
    pdf_path.write_bytes(sample_pdf_bytes)

    result = await PdfExtractor().extract(str(pdf_path))

    assert "Hello World" in result.text
    assert result.source_locations[0].page == 1


@pytest.mark.asyncio
async def test_pdf_extract_from_binary_io(sample_pdf_bytes: bytes) -> None:
    result = await PdfExtractor().extract(io.BytesIO(sample_pdf_bytes))

    assert "Hello World" in result.text
    assert result.source_locations[0].page == 1


@pytest.mark.asyncio
async def test_docx_extract(sample_docx_bytes: bytes) -> None:
    result = await DocxExtractor().extract(sample_docx_bytes)

    assert isinstance(result, ExtractedContent)
    assert "First paragraph hello world" in result.text
    assert "Second paragraph extraction test" in result.text
    assert result.source_type == "docx"
    assert result.metadata["paragraph_count"] == 2
    assert len(result.source_locations) == 2
    assert all(loc.line_start is not None for loc in result.source_locations)
    assert all(loc.line_end is not None for loc in result.source_locations)


@pytest.mark.asyncio
async def test_pptx_extract(sample_pptx_bytes: bytes) -> None:
    result = await PptxExtractor().extract(sample_pptx_bytes)

    assert isinstance(result, ExtractedContent)
    assert "Slide One Title" in result.text
    assert "Slide Two Title" in result.text
    assert result.source_type == "pptx"
    assert result.metadata["slide_count"] == 2
    assert len(result.source_locations) == 2
    assert result.source_locations[0].page == 1
    assert result.source_locations[1].page == 2


@pytest.mark.asyncio
async def test_protocol_conformance() -> None:
    assert isinstance(PdfExtractor(), MaterialExtractor)
    assert isinstance(DocxExtractor(), MaterialExtractor)
    assert isinstance(PptxExtractor(), MaterialExtractor)


def test_source_location_all_fields_optional() -> None:
    loc = SourceLocation()
    assert loc.page is None
    assert loc.line_start is None
    assert loc.timestamp_start_ms is None
    assert loc.bbox is None
